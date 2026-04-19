import json
import subprocess
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.services.slide_matcher import (
    extract_slide_texts,
    frame_number_to_timestamp,
    load_slide_frames,
    match_slides_to_frames,
    render_slides_to_png,
)


# --- Helpers ---

def _make_pptx(tmp_path: Path, slides: list[list[str]]) -> str:
    """Create a .pptx with given slides, each slide is a list of text strings."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for texts in slides:
        slide = prs.slides.add_slide(blank_layout)
        for text in texts:
            txBox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(4), Inches(1))
            txBox.text_frame.text = text
    path = str(tmp_path / "aula.pptx")
    prs.save(path)
    return path


def _make_classifications(tmp_path: Path, entries: list[dict]) -> Path:
    frames_dir = tmp_path / "frames"
    frames_dir.mkdir()
    (frames_dir / "classifications.json").write_text(json.dumps(entries))
    return frames_dir


# --- extract_slide_texts ---

def test_extract_slide_texts_returns_one_entry_per_slide(tmp_path):
    pptx = _make_pptx(tmp_path, [["Slide 1"], ["Slide 2"], ["Slide 3"]])

    result = extract_slide_texts(pptx)

    assert len(result) == 3
    assert all("slide_index" in r and "text" in r for r in result)


def test_extract_slide_texts_correct_index(tmp_path):
    pptx = _make_pptx(tmp_path, [["First"], ["Second"]])

    result = extract_slide_texts(pptx)

    assert result[0]["slide_index"] == 0
    assert result[1]["slide_index"] == 1


def test_extract_slide_texts_concatenates_shapes(tmp_path):
    pptx = _make_pptx(tmp_path, [["Hello", "World"]])

    result = extract_slide_texts(pptx)

    assert "Hello" in result[0]["text"]
    assert "World" in result[0]["text"]


def test_extract_slide_texts_handles_empty_slide(tmp_path):
    pptx = _make_pptx(tmp_path, [[]])

    result = extract_slide_texts(pptx)

    assert result[0]["text"] == ""


# --- load_slide_frames ---

def test_load_slide_frames_filters_to_slide_classification(tmp_path):
    entries = [
        {"frame": "frame_0001.jpg", "frame_path": "/f/frame_0001.jpg", "classification": "slide"},
        {"frame": "frame_0002.jpg", "frame_path": "/f/frame_0002.jpg", "classification": "camera"},
        {"frame": "frame_0003.jpg", "frame_path": "/f/frame_0003.jpg", "classification": "slide"},
    ]
    frames_dir = _make_classifications(tmp_path, entries)

    result = load_slide_frames(str(frames_dir))

    assert len(result) == 2
    assert all(r["classification"] == "slide" for r in result)


def test_load_slide_frames_raises_when_no_classifications_file(tmp_path):
    empty_dir = tmp_path / "noframes"
    empty_dir.mkdir()

    with pytest.raises(FileNotFoundError):
        load_slide_frames(str(empty_dir))


def test_load_slide_frames_returns_empty_when_no_slide_frames(tmp_path):
    entries = [
        {"frame": "frame_0001.jpg", "frame_path": "/f/frame_0001.jpg", "classification": "camera"},
    ]
    frames_dir = _make_classifications(tmp_path, entries)

    result = load_slide_frames(str(frames_dir))

    assert result == []


# --- frame_number_to_timestamp ---

def test_frame_number_to_timestamp_first_frame_is_zero(tmp_path):
    assert frame_number_to_timestamp("frame_0001.jpg", interval=5) == 0


def test_frame_number_to_timestamp_calculates_correctly():
    assert frame_number_to_timestamp("frame_0042.jpg", interval=5) == 205


def test_frame_number_to_timestamp_respects_interval():
    assert frame_number_to_timestamp("frame_0003.jpg", interval=10) == 20


# --- render_slides_to_png ---

def test_render_slides_to_png_calls_libreoffice(tmp_path):
    pptx = str(tmp_path / "aula.pptx")
    Path(pptx).touch()
    out_dir = str(tmp_path / "out")
    Path(out_dir).mkdir()

    with patch("subprocess.run") as mock_run:
        mock_run.return_value = MagicMock(returncode=0)
        render_slides_to_png(pptx, out_dir)

    args = mock_run.call_args[0][0]
    assert "libreoffice" in args
    assert "--headless" in args
    assert pptx in args
    assert out_dir in " ".join(args)


def test_render_slides_to_png_returns_none_on_failure(tmp_path):
    pptx = str(tmp_path / "aula.pptx")
    Path(pptx).touch()
    out_dir = str(tmp_path / "out")
    Path(out_dir).mkdir()

    with patch("subprocess.run", side_effect=subprocess.CalledProcessError(1, "libreoffice")):
        result = render_slides_to_png(pptx, out_dir)

    assert result is None


# --- match_slides_to_frames ---

def test_match_slides_to_frames_returns_best_match_per_slide(tmp_path):
    import torch

    slide_pngs = [str(tmp_path / "slide_0.png")]
    Path(slide_pngs[0]).touch()

    slide_frames = [
        {"frame": "frame_0010.jpg", "frame_path": str(tmp_path / "frame_0010.jpg"), "classification": "slide"},
        {"frame": "frame_0020.jpg", "frame_path": str(tmp_path / "frame_0020.jpg"), "classification": "slide"},
    ]
    for f in slide_frames:
        Path(f["frame_path"]).touch()

    # Mock classifier: first frame has higher similarity
    mock_classifier = MagicMock()
    high_sim = torch.tensor([[1.0, 0.0]])   # normalized: points in same direction as slide
    low_sim  = torch.tensor([[0.0, 1.0]])
    slide_emb = torch.tensor([[1.0, 0.0]])

    def fake_features(pixel_values=None, **_):
        class _Emb:
            def __init__(self, t): self._t = t
            def __truediv__(self, other): return self
            def norm(self, **kw): return torch.tensor(1.0)
        return _Emb(slide_emb)

    mock_classifier.processor = MagicMock(return_value={"pixel_values": torch.zeros(1, 3, 224, 224)})
    # Return slide_emb for slide, high_sim for frame 1, low_sim for frame 2
    embs = [slide_emb, high_sim, low_sim]
    call_count = [0]

    def fake_get_image_features(**kwargs):
        emb = embs[call_count[0] % len(embs)]
        call_count[0] += 1
        norm = emb.norm(dim=-1, keepdim=True)
        return emb / norm if norm > 0 else emb

    mock_classifier.model.get_image_features = fake_get_image_features

    result = match_slides_to_frames(slide_pngs, slide_frames, mock_classifier, interval=5)

    assert len(result) == 1
    assert result[0]["slide_index"] == 0
    assert "start_time" in result[0]
    assert "similarity" in result[0]
