import json
import os
import re
import subprocess
import tempfile
from pathlib import Path


def extract_slide_texts(pptx_path: str) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(pptx_path)
    result = []
    for i, slide in enumerate(prs.slides):
        parts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        text = " ".join(parts).strip()
        result.append({"slide_index": i, "text": text})
    return result


def load_slide_frames(frames_dir: str) -> list[dict]:
    classifications_path = Path(frames_dir) / "classifications.json"
    if not classifications_path.exists():
        raise FileNotFoundError(f"No classifications.json in {frames_dir}")
    entries = json.loads(classifications_path.read_text())
    return [e for e in entries if e.get("classification") == "slide"]


def frame_number_to_timestamp(frame_name: str, interval: int) -> int:
    match = re.search(r"frame_(\d+)", frame_name)
    number = int(match.group(1))
    return (number - 1) * interval


def render_slides_to_png(pptx_path: str, output_dir: str) -> list[str] | None:
    try:
        result = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "png",
                "--outdir",
                output_dir,
                pptx_path,
            ],
            capture_output=True,
        )
        if result.returncode != 0:
            return None
    except subprocess.CalledProcessError:
        return None

    pngs = sorted(Path(output_dir).glob("*.png"))
    return [str(p) for p in pngs]


def _embed_image(path: str, classifier) -> "torch.Tensor":
    from PIL import Image

    image = Image.open(path).convert("RGB")
    inputs = classifier.processor(images=image, return_tensors="pt").to(classifier.device)
    with classifier._torch.no_grad():
        vision_out = classifier.model.vision_model(pixel_values=inputs["pixel_values"])
        emb = classifier.model.visual_projection(vision_out.pooler_output)
    norm = emb.norm(dim=-1, keepdim=True)
    return emb / norm if norm.item() > 0 else emb


def match_slides_to_frames(
    slide_pngs: list[str],
    slide_frames: list[dict],
    classifier,
    interval: int = 5,
) -> list[dict]:
    import torch

    # Pre-embed all frames once, then do a single matmul per slide
    frame_embs = torch.cat([_embed_image(f["frame_path"], classifier) for f in slide_frames])

    results = []
    for slide_index, slide_png in enumerate(slide_pngs):
        slide_emb = _embed_image(slide_png, classifier)
        scores = (slide_emb @ frame_embs.T).squeeze(0)
        best_idx = int(scores.argmax())
        results.append({
            "slide_index": slide_index,
            "frame_path": slide_frames[best_idx]["frame_path"],
            "start_time": frame_number_to_timestamp(slide_frames[best_idx]["frame"], interval),
            "similarity": float(scores[best_idx]),
        })

    return results


def match_slides_to_frames_runpod(
    slide_pngs: list[str],
    slide_frames: list[dict],
    client,
    endpoint_id: str,
    interval: int = 5,
) -> list[dict]:
    import numpy as np
    from app.services.clip_classifier import embed_image_via_runpod

    frame_embs = [embed_image_via_runpod(f["frame_path"], client, endpoint_id) for f in slide_frames]
    frame_matrix = np.array(frame_embs)  # (N, D)
    results = []
    for slide_index, slide_png in enumerate(slide_pngs):
        slide_emb = np.array(embed_image_via_runpod(slide_png, client, endpoint_id))
        scores = frame_matrix @ slide_emb
        best_idx = int(scores.argmax())
        results.append({
            "slide_index": slide_index,
            "frame_path": slide_frames[best_idx]["frame_path"],
            "start_time": frame_number_to_timestamp(slide_frames[best_idx]["frame"], interval),
            "similarity": float(scores[best_idx]),
        })
    return results


def ingest_pptx(
    pptx_path: str,
    video_path: str,
    frames_dir: str,
    interval: int = 5,
    collection_name: str = None,
) -> dict:
    from app.config.settings import Settings
    from app.services.ingestion import (
        _build_point,
        _build_deep_link,
        _extract_class_meta,
        _get_embedding_models,
        _get_ner_pipeline,
        upload_in_batches,
        _get_qdrant_client,
    )
    from app.services.r2_storage import upload_thumbnail, thumbnail_key
    from app.database import get_video_url_by_video_path

    settings = Settings()
    collection_name = collection_name or os.getenv("COLLECTION_NAME", "aulas")

    slide_texts = extract_slide_texts(pptx_path)
    slide_frames = load_slide_frames(frames_dir)

    if not slide_frames:
        return {"ingested": 0, "message": "No slide frames found in classifications"}

    video_url = get_video_url_by_video_path(video_path)

    if settings.use_runpod:
        from app.services.runpod_client import RunPodClient
        clip_client = RunPodClient(settings)
        clip_endpoint = settings.runpod_clip_endpoint_id
        with tempfile.TemporaryDirectory() as tmp_dir:
            pngs = render_slides_to_png(pptx_path, tmp_dir)
            if not pngs:
                return {"ingested": 0, "message": "LibreOffice rendering failed"}
            matches = match_slides_to_frames_runpod(pngs, slide_frames, clip_client, clip_endpoint, interval=interval)
            thumb_urls: dict[int, str | None] = {}  # skip thumbnail upload in RunPod mode for now
        embedding_arg = settings
        ner_arg = settings
    else:
        from app.services.clip_classifier import get_classifier
        classifier = get_classifier()
        with tempfile.TemporaryDirectory() as tmp_dir:
            pngs = render_slides_to_png(pptx_path, tmp_dir)
            if not pngs:
                return {"ingested": 0, "message": "LibreOffice rendering failed"}

            matches = match_slides_to_frames(pngs, slide_frames, classifier, interval=interval)

            # Upload thumbnails while temp PNGs still exist
            thumb_urls = {}
            for match in matches:
                idx = match["slide_index"]
                if idx < len(pngs):
                    key = thumbnail_key(video_path, idx)
                    thumb_urls[idx] = upload_thumbnail(pngs[idx], key)

        embedding_arg = _get_embedding_models()
        ner_arg = _get_ner_pipeline()

    class_meta = _extract_class_meta(video_path)
    client = _get_qdrant_client()

    slide_text_map = {s["slide_index"]: s["text"] for s in slide_texts}
    points = []
    for match in matches:
        idx = match["slide_index"]
        text = slide_text_map.get(idx, "")
        if not text.strip():
            continue
        payload = {
            "source_type": "slide",
            "file_path": pptx_path,
            "video_url": _build_deep_link(video_url, match["start_time"]),
            "slide_index": idx,
            "start_time": match["start_time"],
            **class_meta,
        }
        if thumb_urls.get(idx):
            payload["slide_thumb"] = thumb_urls[idx]
        points.append(_build_point(text, embedding_arg, ner_arg, payload))

    if points:
        upload_in_batches(client, collection_name, points)

    return {"ingested": len(points)}
